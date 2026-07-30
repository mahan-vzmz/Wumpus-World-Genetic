from __future__ import annotations

import csv
import html
import json
import shutil
import subprocess
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
REPORT_DIR = DOCS / "final_report"
PRESENTATION_DIR = DOCS / "presentation"
RESULTS = ROOT / "results" / "final"


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle))


def copy_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    sources = [
        RESULTS / "success_rate.png",
        RESULTS / "average_score.png",
        RESULTS / "average_steps_success.png",
        RESULTS / "remaining_health.png",
        RESULTS / "runtime.png",
        RESULTS / "success_by_difficulty.png",
        RESULTS / "failure_reasons.png",
        ROOT / "results" / "genetic_fitness.png",
    ]
    for source in sources:
        if not source.exists():
            raise FileNotFoundError(
                f"Missing artifact source: {source}. Run train_genetic.py and experiment.py first."
            )
        shutil.copy2(source, ASSETS / source.name)


def build_report(info: dict[str, str], summary: list[dict[str, str]]) -> Path:
    from weasyprint import HTML

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    rows = "".join(
        f"<tr><td>{html.escape(row['agent'])}</td>"
        f"<td>{row['success_rate']}%</td>"
        f"<td>{row['average_score_all']}</td>"
        f"<td>{row['average_steps_all']}</td>"
        f"<td>{row['average_steps_success']}</td>"
        f"<td>{row['average_remaining_health_all']}</td>"
        f"<td>{row['average_pit_entries']}</td>"
        f"<td>{row['wumpus_deaths']}</td></tr>"
        for row in summary
    )

    html_text = f"""<!doctype html>
<html lang="fa" dir="rtl">
<head>
<meta charset="utf-8">
<style>
@page {{ size: A4; margin: 18mm 16mm 18mm 16mm; @bottom-center {{ content: counter(page); font-size: 9pt; }} }}
body {{ font-family: 'Noto Sans Arabic', 'DejaVu Sans', sans-serif; direction: rtl; color:#172033; line-height:1.75; font-size:10.5pt; }}
h1,h2,h3 {{ color:#123b5d; page-break-after: avoid; }}
h1 {{ font-size:23pt; text-align:center; margin-top:45mm; }}
h2 {{ font-size:16pt; border-bottom:1px solid #cbd7e1; padding-bottom:3px; margin-top:18px; }}
h3 {{ font-size:12.5pt; }}
p {{ text-align:justify; }}
.cover {{ page-break-after: always; text-align:center; }}
.meta {{ margin:28mm auto 0; width:82%; border:1px solid #b8c6d2; border-radius:8px; padding:18px; background:#f6f9fb; }}
.meta p {{ text-align:center; margin:8px; }}
table {{ width:100%; border-collapse:collapse; margin:12px 0; font-size:9pt; direction:ltr; }}
th,td {{ border:1px solid #9fb1c1; padding:6px; text-align:center; }}
th {{ background:#e8f0f6; color:#123b5d; }}
figure {{ page-break-inside:avoid; margin:14px auto; text-align:center; }}
figure img {{ max-width:95%; max-height:95mm; }}
figcaption {{ font-size:9pt; color:#566; margin-top:4px; }}
.callout {{ background:#eef6fb; border-right:4px solid #1f76b4; padding:9px 12px; margin:12px 0; }}
.code {{ direction:ltr; text-align:left; font-family:monospace; background:#f2f4f6; padding:8px; }}
ul {{ margin-right:20px; }}
.small {{ font-size:9pt; color:#526273; }}
</style>
</head>
<body>
<section class="cover">
<h1>گزارش نهایی پروژه Wumpus World</h1>
<h2 style="border:0;text-align:center">نسخه 8 - مقایسه A-Star، عامل قاعده‌محور و عامل ژنتیکی ترکیبی</h2>
<div class="meta">
<p><b>نام دانشجو:</b> {html.escape(info['student_name'])}</p>
<p><b>شماره دانشجویی:</b> {html.escape(info['student_id'])}</p>
<p><b>درس:</b> {html.escape(info['course_name'])}</p>
<p><b>استاد:</b> {html.escape(info['instructor_name'])}</p>
<p><b>دانشگاه:</b> {html.escape(info['university_name'])}</p>
<p><b>تاریخ:</b> {html.escape(info['submission_date'])}</p>
</div>
</section>

<h2>چکیده</h2>
<p>در این پروژه محیط Wumpus World روی گرید 8×8 پیاده‌سازی شد و سه روش متفاوت روی یک محیط و مجموعه معیار مشترک مقایسه شدند. A-Star به کل نقشه دسترسی دارد و نقش Oracle را ایفا می‌کند. عامل قاعده‌محور و عامل ژنتیکی ترکیبی فقط از ادراک‌های محلی، حافظه و مختصات عمومی خروج استفاده می‌کنند. وزن‌های روش ژنتیکی روی 12 نقشه آموزش تکامل یافتند و ارزیابی نهایی روی 30 نقشه تست جداگانه و 90 اپیزود انجام شد.</p>
<div class="callout">نتیجه اصلی: A-Star به 100٪، Rule-Based به 90٪ و Hybrid Genetic به 83.33٪ موفقیت رسید. در میان عامل‌های آنلاین، Rule-Based مطمئن‌تر و Hybrid Genetic در اپیزودهای موفق کوتاه‌مسیرتر بود.</div>

<h2>۱. تعریف مسئله و قوانین</h2>
<p>عامل از خانه (1,1) شروع می‌کند، باید حداقل یک طلا جمع‌آوری کند و پیش از تمام‌شدن جان به خروج برسد. هر تلاش برای حرکت یک واحد جان کم می‌کند. دیوار حرکت را مسدود می‌کند، چاه پس از هزینه حرکت جان را نصف می‌کند، غول مرگ فوری ایجاد می‌کند و خروج بدون طلا شکست است.</p>
<ul><li>Breeze: وجود چاه در همسایگی چهارجهته</li><li>Stench: وجود غول در همسایگی چهارجهته</li><li>Pit here: تشخیص چاه پس از زنده‌ماندن و ورود</li><li>حرکت‌های معتبر، جان، داشتن طلا و مختصات خروج</li></ul>

<h2>۲. معماری</h2>
<p>parser، محیط، عامل‌ها، پایگاه دانش، آموزش ژنتیک، مولد نقشه و benchmark در ماژول‌های مستقل قرار گرفته‌اند. تمام عامل‌ها از تابع اجرای مشترک استفاده می‌کنند. نسخه 8 باگ پایان max_steps، برنامه‌ریزی دوباره A-Star، تشخیص نادرست چاه بازدیدشده، fallback بی‌صدای وزن‌های ژنتیکی و اعتبارسنجی ناقص نقشه را اصلاح می‌کند.</p>

<h2>۳. روش A-Star</h2>
<p>حالت جست‌وجو شامل موقعیت، جان باقی‌مانده و داشتن طلاست. دیوار و غول از فضای حالت حذف می‌شوند. ورود به چاه در صورت زنده‌ماندن مجاز است، اما کاهش واقعی جان و جریمه چاه در هزینه مسیر لحاظ می‌شود. heuristic فاصله منهتن یک lower bound معتبر است.</p>

<h2>۴. عامل قاعده‌محور</h2>
<p>این عامل از Breeze و Stench پایگاه دانش می‌سازد. نبود ادراک خطر، همسایه‌ها را از همان خطر مبرا می‌کند. clause تک‌عضوی محل خطر قطعی را مشخص می‌کند. عامل خانه امن بازدیدنشده را ترجیح می‌دهد، در بن‌بست backtracking می‌کند و در نبود گزینه امن، کم‌خطرترین frontier را انتخاب می‌کند.</p>

<h2>۵. عامل ژنتیکی ترکیبی</h2>
<p>روش سوم یک عامل ترکیبی است. پایگاه دانش محلی شواهد خطر را فراهم می‌کند؛ یک سیاست خطی با 10 وزن تکامل‌یافته، حرکت‌های مرحله اکتشاف را امتیازدهی می‌کند؛ و پس از گرفتن طلا، عامل از کوتاه‌ترین مسیر شناخته‌شده امن به خروج بازمی‌گردد.</p>
<p class="code">score(action) = Σ weight_i × feature_i(action)</p>

<h2>۶. آموزش الگوریتم ژنتیک</h2>
<ul><li>12 نقشه آموزش جدا</li><li>Population = 24</li><li>Maximum generations = 24</li><li>Elitism = 2</li><li>Tournament size = 3</li><li>Mutation rate = 0.10</li><li>Seed = 17</li><li>Best fitness = 1840.67</li></ul>
<figure><img src="../assets/genetic_fitness.png"><figcaption>روند بهترین و میانگین Fitness در طول آموزش</figcaption></figure>

<h2>۷. طراحی آزمایش</h2>
<p>30 نقشه تست دیده‌نشده شامل 10 نقشه آسان، 10 متوسط و 10 سخت با seed ثابت تولید شدند. خروج، محل طلا و طول مسیرها متنوع است. جان اولیه همه سطوح برابر 120 است تا امتیاز بین دشواری‌ها قابل مقایسه باشد. زمان اجرا median سه اجرای کامل است و تعداد حرکت موفق جداگانه گزارش می‌شود.</p>

<h2>۸. نتایج کلی</h2>
<table><thead><tr><th>Agent</th><th>Success</th><th>Score all</th><th>Steps all</th><th>Steps success</th><th>Health</th><th>Pit entries</th><th>Wumpus deaths</th></tr></thead><tbody>{rows}</tbody></table>
<figure><img src="../assets/success_rate.png"><figcaption>نرخ موفقیت سه روش</figcaption></figure>
<figure><img src="../assets/average_steps_success.png"><figcaption>میانگین حرکت فقط در اپیزودهای موفق</figcaption></figure>

<h2>۹. تحلیل مقایسه‌ای</h2>
<p>A-Star به دلیل اطلاعات کامل و وجود حداقل یک مسیر امن، کران بالای عملکرد است. Rule-Based در میان عامل‌های آنلاین نرخ موفقیت بالاتری دارد و محافظه‌کارتر است. Hybrid Genetic در اپیزودهای موفق حرکت کمتری مصرف می‌کند و امتیاز موفقیت بالاتری دارد، اما ریسک مرگ با غول بیشتر است. بنابراین نتیجه علمی، برتری مطلق یک روش نیست؛ بلکه تفاوت در trade-off میان اطمینان، توضیح‌پذیری و سرعت موفقیت است.</p>
<figure><img src="../assets/success_by_difficulty.png"><figcaption>نرخ موفقیت بر اساس سطح دشواری</figcaption></figure>
<figure><img src="../assets/failure_reasons.png"><figcaption>دلایل شکست</figcaption></figure>

<h2>۱۰. محدودیت‌ها</h2>
<ul><li>A-Star سطح اطلاعات متفاوتی دارد.</li><li>نتایج برای seed و مجموعه تست ثبت‌شده معتبرند.</li><li>عامل ژنتیکی تضمین بهینگی یا موفقیت ندارد.</li><li>زمان اجرا به سخت‌افزار وابسته است.</li><li>برای استنباط آماری قوی‌تر، چند seed و confidence interval لازم است.</li></ul>

<h2>۱۱. نتیجه‌گیری</h2>
<p>نسخه 8 یک pipeline قابل‌بازتولید از تولید نقشه و آموزش تا تست، ارزیابی، گزارش و ارائه فراهم می‌کند. A-Star بهترین عملکرد را در محیط کاملاً شناخته‌شده دارد. در محیط ناشناخته، Rule-Based مطمئن‌تر و توضیح‌پذیرتر است، در حالی که Hybrid Genetic در موفقیت‌ها کوتاه‌مسیرتر اما ریسک‌پذیرتر عمل می‌کند.</p>

<h2>۱۲. منابع</h2>
<ol><li>Russell, S. J., & Norvig, P. Artificial Intelligence: A Modern Approach.</li><li>Hart, P. E., Nilsson, N. J., & Raphael, B. (1968). A Formal Basis for the Heuristic Determination of Minimum Cost Paths.</li><li>Holland, J. H. (1975). Adaptation in Natural and Artificial Systems.</li><li>Goldberg, D. E. (1989). Genetic Algorithms in Search, Optimization, and Machine Learning.</li></ol>
</body></html>"""
    html_path = REPORT_DIR / "final_report.html"
    pdf_path = REPORT_DIR / "final_report.pdf"
    html_path.write_text(html_text, encoding="utf-8")
    HTML(filename=str(html_path), base_url=str(REPORT_DIR)).write_pdf(str(pdf_path))
    return pdf_path


def add_textbox(slide, text: str, left: float, top: float, width: float, height: float, *, size: int = 20, bold: bool = False, color: str = "25364A", align: str = "right") -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = {"right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT}[align]
    paragraph.font.name = "Noto Sans Arabic"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(color)


def build_presentation(info: dict[str, str], summary: list[dict[str, str]]) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    navy = RGBColor(15, 54, 88)
    blue = RGBColor(31, 118, 180)
    light = RGBColor(244, 248, 251)
    gray = RGBColor(92, 112, 130)

    def base_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        add_textbox(slide, title, 0.65, 0.28, 12.0, 0.55, size=29, bold=True, color="0F3658")
        if subtitle:
            add_textbox(slide, subtitle, 0.65, 0.88, 12.0, 0.35, size=13, color="5C7082")
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.32), Inches(12.0), Inches(0.025))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(214, 225, 234); line.line.fill.background()
        add_textbox(slide, f"Wumpus World v8  |  {info['student_name']}", 0.65, 7.08, 12.0, 0.22, size=9, color="6B7D8C", align="center")
        return slide

    def bullets(slide, items: list[str], top: float = 1.65, size: int = 21):
        box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.6), Inches(5.1))
        frame = box.text_frame; frame.clear(); frame.word_wrap = True
        for index, item in enumerate(items):
            p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            p.text = item; p.alignment = PP_ALIGN.RIGHT
            p.font.name = "Noto Sans Arabic"; p.font.size = Pt(size); p.font.color.rgb = RGBColor(37,54,74)
            p.space_after = Pt(12)
            p.level = 0

    # 1 cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = navy
    add_textbox(slide, "Wumpus World", 0.8, 1.15, 11.7, 0.8, size=42, bold=True, color="FFFFFF", align="center")
    add_textbox(slide, "نسخه 8 - مقایسه A-Star، عامل قاعده‌محور و عامل ژنتیکی ترکیبی", 1.0, 2.1, 11.3, 0.7, size=24, color="DCEAF5", align="center")
    add_textbox(slide, f"{info['student_name']}  |  {info['course_name']}  |  {info['submission_date']}", 1.0, 5.65, 11.3, 0.45, size=17, color="FFFFFF", align="center")

    slide = base_slide("تعریف مسئله", "هدف عامل: طلا را بگیرد و زنده به خروج برسد")
    bullets(slide, ["گرید 8×8 با دیوار، چاه، غول، طلا و خروج", "هر حرکت یک واحد جان کم می‌کند", "چاه جان را نصف و غول عامل را نابود می‌کند", "Breeze و Stench تنها ادراک خطر برای عامل‌های آنلاین هستند"])

    slide = base_slide("معماری مشترک", "هر سه روش روی محیط و قرارداد اجرای یکسان")
    bullets(slide, ["Parser سخت‌گیرانه و محیط deterministic", "رابط مشترک BaseAgent و تابع run_episode", "پایگاه دانش محلی برای دو عامل آنلاین", "Pipeline کامل: تولید نقشه → آموزش → تست → CSV → گزارش"])

    slide = base_slide("روش اول: A-Star", "Oracle با اطلاعات کامل نقشه")
    bullets(slide, ["حالت: موقعیت، جان و داشتن طلا", "دیوار و غول غیرقابل عبور", "چاه با کاهش واقعی جان و جریمه وارد cost می‌شود", "Heuristic: فاصله منهتن تا طلا و سپس خروج", "Baseline بالادستی؛ مقایسه مستقیم با عامل‌های آنلاین منصفانه نیست"])

    slide = base_slide("روش دوم: Rule-Based", "استنتاج محلی، خانه امن و backtracking")
    bullets(slide, ["No Breeze ⇒ همسایه‌ها چاه ندارند", "No Stench ⇒ همسایه‌ها غول ندارند", "Clause تک‌عضوی ⇒ خطر قطعی", "اولویت با safe frontier، سپس backtracking", "در نبود گزینه امن، کم‌خطرترین frontier انتخاب می‌شود"])

    slide = base_slide("روش سوم: Hybrid Genetic", "وزن‌های GA برای اکتشاف + پایگاه دانش + بازگشت امن")
    bullets(slide, ["10 ویژگی و 10 وزن حقیقی", "score(action) = Σ(weight × feature)", "وزن‌ها روی 12 نقشه آموزش تکامل یافته‌اند", "پس از طلا، کوتاه‌ترین مسیر شناخته‌شده امن استفاده می‌شود", "روش ترکیبی است؛ نه یک عامل کاملاً مستقل از قواعد"])

    slide = base_slide("آموزش الگوریتم ژنتیک", "Population=24 | Generations=24 | Seed=17")
    slide.shapes.add_picture(str(ASSETS / "genetic_fitness.png"), Inches(1.15), Inches(1.55), width=Inches(7.6), height=Inches(4.75))
    add_textbox(slide, "Best fitness\n1840.67\n\nTraining maps\n12\n\nTraining success\n100%", 9.2, 1.8, 3.1, 4.1, size=22, bold=True, color="0F3658", align="center")

    slide = base_slide("طراحی آزمایش نهایی", "30 نقشه دیده‌نشده و 90 اپیزود")
    bullets(slide, ["10 آسان، 10 متوسط، 10 سخت", "خروج، طلا و طول مسیرها متنوع", "جان اولیه همه سطوح برابر 120", "زمان: median سه اجرای کامل", "حرکت اپیزودهای موفق جداگانه گزارش می‌شود"])

    slide = base_slide("نتیجه کلی", "A-Star کران بالا؛ مقایسه اصلی بین دو عامل آنلاین")
    slide.shapes.add_picture(str(ASSETS / "success_rate.png"), Inches(0.8), Inches(1.55), width=Inches(7.2), height=Inches(4.65))
    data = {row["agent"]: row for row in summary}
    text = (
        f"A-Star: {data['astar']['success_rate']}%\n"
        f"Rule-Based: {data['rule']['success_rate']}%\n"
        f"Hybrid Genetic: {data['genetic']['success_rate']}%"
    )
    add_textbox(slide, text, 8.35, 2.0, 4.1, 2.8, size=24, bold=True, color="0F3658", align="center")

    slide = base_slide("تحلیل عامل‌های آنلاین", "اطمینان بیشتر در برابر مسیر موفق کوتاه‌تر")
    slide.shapes.add_picture(str(ASSETS / "average_steps_success.png"), Inches(0.75), Inches(1.55), width=Inches(7.2), height=Inches(4.65))
    bullets(slide, ["Rule-Based: موفقیت 90٪", "Hybrid Genetic: موفقیت 83.33٪", "حرکت موفق Rule-Based: 32.30", "حرکت موفق Genetic: 24.60", "Genetic سریع‌تر اما در برابر غول ریسک‌پذیرتر است"], top=1.72, size=18)
    # move bullet box to the right by adjusting the last shape
    last = slide.shapes[-1]; last.left = Inches(8.15); last.width = Inches(4.4); last.height = Inches(4.7)

    slide = base_slide("اصلاحات کلیدی نسخه 8", "مشکلات منطقی، آزمایشی و مستنداتی برطرف شدند")
    bullets(slide, ["ثبت صحیح max_steps و علت پایان", "تشخیص چاه بازدیدشده و جلوگیری از safe اشتباه", "حذف برنامه‌ریزی دوباره A-Star در timing", "جان اولیه ثابت و خروج‌های متنوع", "معیار حرکت موفق و runtime تکرارشده", "44 تست، CI، MIT License و artifacts قابل‌بازتولید"])

    slide = base_slide("محدودیت‌ها", "تفسیر دقیق و قابل دفاع")
    bullets(slide, ["A-Star سطح اطلاعات متفاوت دارد", "نتایج برای seed ثبت‌شده معتبرند", "GA تضمین موفقیت یا بهینگی ندارد", "Runtime به سخت‌افزار وابسته است", "چند seed و فاصله اطمینان، ادامه علمی مناسب پروژه است"])

    slide = base_slide("نتیجه‌گیری", "انتخاب روش به سطح اطلاعات و اولویت پروژه وابسته است")
    bullets(slide, ["نقشه کامل: A-Star بهترین baseline", "محیط ناشناخته و نیاز به اطمینان: Rule-Based", "کوتاهی مسیر موفق با پذیرش ریسک بیشتر: Hybrid Genetic", "نسخه 8 شامل کد، 44 تست، داده، وزن، نتایج، گزارش و ارائه است"])

    output = PRESENTATION_DIR / "wumpus_world_presentation.pptx"
    prs.save(output)
    return output


def export_presentation_pdf(pptx_path: Path) -> Path | None:
    executable = shutil.which("libreoffice")
    if executable is None:
        return None
    subprocess.run(
        [
            executable,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(PRESENTATION_DIR),
            str(pptx_path),
        ],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    return PRESENTATION_DIR / (pptx_path.stem + ".pdf")


def main() -> None:
    info = json.loads((ROOT / "project_info.json").read_text(encoding="utf-8"))
    summary = read_csv(RESULTS / "summary_results.csv")
    copy_assets()
    report = build_report(info, summary)
    presentation = build_presentation(info, summary)
    presentation_pdf = export_presentation_pdf(presentation)
    print(f"report={report.relative_to(ROOT)}")
    print(f"presentation={presentation.relative_to(ROOT)}")
    if presentation_pdf and presentation_pdf.exists():
        print(f"presentation_pdf={presentation_pdf.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
